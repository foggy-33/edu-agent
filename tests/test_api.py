import base64
from io import BytesIO

from docx import Document
from fastapi.testclient import TestClient
from pptx import Presentation

from app.auth.service import AuthError, AuthService
from app.core.config import get_settings
from app.main import app


client = TestClient(app)


def test_health() -> None:
    response = client.get("/health")
    assert response.status_code == 200
    assert response.json()["status"] == "ok"


def test_register_login_me_and_logout() -> None:
    username = "pytest_auth_user"
    register = client.post(
        "/api/auth/register",
        json={"username": username, "display_name": "测试用户", "password": "test-pass-123"},
    )
    if register.status_code == 400:
        login = client.post("/api/auth/login", json={"username": username, "password": "test-pass-123"})
        assert login.status_code == 200
        auth = login.json()
    else:
        assert register.status_code == 201
        auth = register.json()

    headers = {"Authorization": f"Bearer {auth['token']}"}
    me = client.get("/api/auth/me", headers=headers)
    assert me.status_code == 200
    assert me.json()["user"]["username"] == username

    logout = client.post("/api/auth/logout", headers=headers)
    assert logout.status_code == 204
    assert client.get("/api/auth/me", headers=headers).status_code == 401


def test_regular_user_cannot_access_admin_users() -> None:
    username = "pytest_regular_admin_guard"
    registered = client.post(
        "/api/auth/register",
        json={"username": username, "display_name": "普通用户", "password": "test-pass-123"},
    )
    auth = registered.json() if registered.status_code == 201 else client.post(
        "/api/auth/login",
        json={"username": username, "password": "test-pass-123"},
    ).json()
    response = client.get(
        "/api/auth/admin/users",
        headers={"Authorization": f"Bearer {auth['token']}"},
    )
    assert response.status_code == 403


def test_admin_bootstrap_and_safe_user_listing(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("USER_DATA_FILE", str(tmp_path / "users.json"))
    monkeypatch.setenv("ADMIN_USERNAME", "local_admin")
    monkeypatch.setenv("ADMIN_PASSWORD", "secure-admin-pass")
    monkeypatch.setenv("ADMIN_DISPLAY_NAME", "测试管理员")
    get_settings.cache_clear()
    try:
        service = AuthService()
        admin = service.login("local_admin", "secure-admin-pass")
        assert admin["user"]["role"] == "admin"
        service.register("student_001", "学生一", "student-pass-123")
        listing = service.list_users()
        assert listing["total"] == 2
        assert {item["username"] for item in listing["users"]} == {"local_admin", "student_001"}
        assert all("password_hash" not in item and "sessions" not in item for item in listing["users"])
        try:
            service.register("local_admin", "冒用管理员", "another-pass-123")
        except AuthError as exc:
            assert "管理员保留账号" in str(exc)
        else:
            raise AssertionError("管理员保留账号不应允许公开注册")
    finally:
        get_settings.cache_clear()


def test_uploaded_avatar_can_be_loaded_from_returned_url() -> None:
    username = "pytest_avatar_user"
    registered = client.post(
        "/api/auth/register",
        json={"username": username, "display_name": "头像测试", "password": "test-pass-123"},
    )
    if registered.status_code == 400:
        login = client.post("/api/auth/login", json={"username": username, "password": "test-pass-123"})
        assert login.status_code == 200
        auth = login.json()
    else:
        assert registered.status_code == 201
        auth = registered.json()

    # Valid 1x1 PNG. The upload response URL must be readable through /api/avatars/.
    image = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
    )
    response = client.post(
        "/api/auth/avatar",
        headers={"Authorization": f"Bearer {auth['token']}"},
        files={"file": ("avatar.png", image, "image/png")},
    )
    assert response.status_code == 200
    avatar_url = response.json()["user"]["avatar"]
    assert avatar_url.startswith("/api/avatars/")

    loaded = client.get(avatar_url)
    assert loaded.status_code == 200
    assert loaded.headers["content-type"] == "image/png"
    assert loaded.content == image


def test_analyze() -> None:
    response = client.post(
        "/api/analyze",
        json={
            "user_id": "demo_user_001",
            "course": "数据库系统",
            "message": "我是计算机专业大二学生，对函数依赖、候选码和范式判断不太会，希望通过例题准备考试。",
        },
    )
    assert response.status_code == 200
    data = response.json()
    assert "profile" in data
    assert "函数依赖" in data["profile"]["weak_points"]


def test_generate() -> None:
    response = client.post(
        "/api/generate",
        json={
            "user_id": "demo_user_001",
            "course": "数据库系统",
            "message": "我是计算机专业大二学生，正在学习数据库系统。我对函数依赖、候选码和范式判断不太会，希望通过例题和步骤化讲解准备考试。",
        },
    )
    assert response.status_code == 200
    data = response.json()
    assert data["resources"]["mindmap"].startswith("mindmap")
    assert len(data["resources"]["quiz"]) >= 5
    assert "CREATE TABLE" in data["resources"]["practice_case"]


def test_export_real_office_files() -> None:
    content = "# 数据库系统\n\n## 核心概念\n- 关系模型\n- 主键与外键\n\n## 复习建议\n1. 整理概念\n2. 完成练习"
    for office_format, expected_type in (
        ("pptx", "presentationml.presentation"),
        ("docx", "wordprocessingml.document"),
    ):
        response = client.post(
            "/api/exports/office",
            json={
                "title": "数据库系统复习",
                "subtitle": "智学 AI 生成",
                "content": content,
                "format": office_format,
            },
        )
        assert response.status_code == 200
        assert expected_type in response.headers["content-type"]
        assert f".{office_format}" in response.headers["content-disposition"]
        assert response.content.startswith(b"PK")
        assert len(response.content) > 1000
        if office_format == "pptx":
            deck = Presentation(BytesIO(response.content))
            assert len(deck.slides) >= 5
            visible_text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text_frame"))
            assert "数据库系统复习" in visible_text
            assert "这次讲解将解决什么" in visible_text
        else:
            document = Document(BytesIO(response.content))
            assert any("数据库系统复习" in paragraph.text for paragraph in document.paragraphs)


def test_collaborative_learning_generate() -> None:
    response = client.post(
        "/api/learning/generate",
        json={
            "major": "计算机科学与技术",
            "course": "操作系统",
            "chapter": "进程调度",
            "weakness": "不会区分 FCFS、SJF 和时间片轮转",
            "goal": "期末复习",
            "resourceTypes": ["lecture", "mindmap", "exercise", "reading"],
            "api_key": "",
            "base_url": "https://api.siliconflow.cn/v1",
            "model": "Pro/deepseek-ai/DeepSeek-V3.2",
        },
    )
    assert response.status_code == 200
    data = response.json()
    assert data["lectureDoc"].startswith("#")
    assert data["mindmap"].startswith("mindmap")
    assert "选择题" in data["exercises"]
    assert data["exerciseItems"][0]["question"]
    assert data["exerciseItems"][0]["answer"]
    assert data["exerciseItems"][0]["explanation"]
    assert data["reading"]
    assert "质量审核" in data["review"]
    assert len(data["agentTrace"]) == 8
    assert data["agentTrace"][0]["agent"] == "学情分析 Agent"
    assert data["agentTrace"][-1]["agent"] == "资源整合 Agent"


def test_collaborative_learning_uses_user_api_settings(monkeypatch) -> None:
    captured = []

    def fake_call_llm(prompt: str, **kwargs) -> str:
        captured.append(kwargs)
        return "# 用户模型生成内容"

    monkeypatch.setattr("app.learning.agents.call_llm", fake_call_llm)
    response = client.post(
        "/api/learning/generate",
        json={
            "major": "计算机科学与技术",
            "course": "操作系统",
            "chapter": "进程调度",
            "weakness": "调度算法对比",
            "goal": "期末复习",
            "resourceTypes": ["lecture"],
            "api_key": "user-provided-key",
            "base_url": "https://example.test/v1",
            "model": "user-model",
        },
    )
    assert response.status_code == 200
    assert captured
    assert all(item["api_key"] == "user-provided-key" for item in captured)
    assert all(item["base_url"] == "https://example.test/v1" for item in captured)
    assert all(item["model"] == "user-model" for item in captured)


def test_pdf_resource_upload_and_generate(monkeypatch) -> None:
    user_id = "pytest_pdf_resource_user"
    monkeypatch.setattr(
        "app.api.routes.resource_service._extract_pages",
        lambda _: [
            {"page": 1, "text": "操作系统进程调度包含先来先服务、短作业优先和时间片轮转。"},
            {"page": 2, "text": "时间片轮转需要在响应速度和上下文切换开销之间权衡。"},
        ],
    )
    uploaded = client.post(
        "/api/resources/upload",
        data={"user_id": user_id},
        files={"file": ("scheduling.pdf", b"%PDF-test", "application/pdf")},
    )
    assert uploaded.status_code == 201
    resource = uploaded.json()["resource"]
    assert resource["name"] == "scheduling.pdf"
    assert resource["page_count"] == 2

    listed = client.get("/api/resources", params={"user_id": user_id})
    assert listed.status_code == 200
    assert any(item["id"] == resource["id"] for item in listed.json()["resources"])

    generated = client.post(
        "/api/learning/generate",
        json={
            "user_id": user_id,
            "major": "未指定",
            "course": "自定义学习主题",
            "chapter": "用户当前问题",
            "weakness": "根据上传资料生成复习内容",
            "goal": "理解并掌握相关知识",
            "resourceTypes": ["lecture", "mindmap", "exercise"],
            "fileIds": [resource["id"]],
            "api_key": "",
            "base_url": "https://api.siliconflow.cn/v1",
            "model": "Pro/deepseek-ai/DeepSeek-V3.2",
        },
    )
    assert generated.status_code == 200
    data = generated.json()
    assert data["sources"][0]["name"] == "scheduling.pdf"
    assert "进程调度" in data["lectureDoc"]
    assert "进程调度" in data["mindmap"]

    deleted = client.delete(f"/api/resources/{resource['id']}", params={"user_id": user_id})
    assert deleted.status_code == 204


def test_dynamic_profile_chat_falls_back_without_api_key() -> None:
    response = client.post(
        "/api/profiles/chat",
        json={
            "user_id": "pytest_profile_user",
            "course": "数据库系统",
            "message": "我是计算机专业大三学生，准备期末考试，范式判断容易出错，喜欢通过例题学习。",
            "api_key": "",
            "base_url": "https://api.siliconflow.cn/v1",
            "model": "Pro/deepseek-ai/DeepSeek-V3.2",
        },
    )
    assert response.status_code == 200
    data = response.json()
    assert data["provider"] == "rule-fallback"
    assert data["profile"]["version"] >= 1
    assert len(data["profile"]["dimension_catalog"]) >= 6
    assert len(data["profile"]["dimensions"]) >= 6

    saved = client.get("/api/profiles/pytest_profile_user")
    assert saved.status_code == 200
    assert saved.json()["profile"]["version"] >= 1


def test_evaluation_updates_dynamic_profile() -> None:
    response = client.post(
        "/api/evaluate",
        json={
            "user_id": "pytest_evaluation_profile_user",
            "course": "数据库系统",
            "answers": [
                {"question": "判断范式", "student_answer": "2NF", "is_correct": False, "topic": "范式判断"},
                {"question": "SQL 查询", "student_answer": "SELECT", "is_correct": True, "topic": "SQL"},
            ],
        },
    )
    assert response.status_code == 200
    profile = response.json()["dynamic_profile"]
    assert profile["version"] >= 1
    assert "易错点" in profile["dimensions"]
    assert "知识基础" in profile["dimensions"]


def test_profile_interview_asks_course_question_without_api_key() -> None:
    response = client.post(
        "/api/profiles/interview/next",
        json={
            "user_id": "pytest_interview_user",
            "course": "数据结构",
            "api_key": "",
            "base_url": "https://api.siliconflow.cn/v1",
            "model": "Pro/deepseek-ai/DeepSeek-V3.2",
        },
    )
    assert response.status_code == 200
    data = response.json()
    assert data["provider"] == "rule-fallback"
    assert data["question"]
    assert len(data["profile"]["dimension_catalog"]) >= 6


def test_dynamic_profiles_are_stored_by_subject() -> None:
    user_id = "pytest_subject_profile_user"
    for course, message in [
        ("数据库系统", "我在准备数据库期末考试，范式判断容易出错，喜欢通过例题学习。"),
        ("操作系统", "我想深入学习操作系统，进程调度比较薄弱，喜欢视频和练习题。"),
    ]:
        response = client.post(
            "/api/profiles/chat",
            json={
                "user_id": user_id,
                "course": course,
                "message": message,
                "api_key": "",
                "base_url": "https://api.siliconflow.cn/v1",
                "model": "Pro/deepseek-ai/DeepSeek-V3.2",
            },
        )
        assert response.status_code == 200
        assert response.json()["profile"]["course"] == course

    listed = client.get(f"/api/profiles/{user_id}/subjects")
    assert listed.status_code == 200
    assert {item["course"] for item in listed.json()["profiles"]} >= {"数据库系统", "操作系统"}

    database_profile = client.get(f"/api/profiles/{user_id}", params={"course": "数据库系统"}).json()["profile"]
    os_profile = client.get(f"/api/profiles/{user_id}", params={"course": "操作系统"}).json()["profile"]
    assert "范式判断" in database_profile["dimensions"]["易错点"]["value"]
    assert "进程调度" in os_profile["dimensions"]["易错点"]["value"]
    assert database_profile["llm_context"]["schema_version"] == "subject-profile-v1"
    assert len(database_profile["radar_metrics"]) == 6
    assert "画像可信" not in database_profile["radar_metrics"]
    assert set(database_profile["radar_metrics"]) == set(database_profile["radar_summaries"])
