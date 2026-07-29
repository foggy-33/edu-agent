from __future__ import annotations

import hashlib
import re
from dataclasses import dataclass
from io import BytesIO
from typing import Iterable

from PIL import Image, ImageDraw
from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt


@dataclass(frozen=True)
class OfficeTheme:
    name: str
    accent: str
    accent_2: str
    ink: str
    muted: str
    soft: str
    canvas: str


THEMES = (
    OfficeTheme("深海科技", "176B87", "64CCC5", "102A43", "627D98", "EAF8F8", "F7FBFC"),
    OfficeTheme("现代学院", "1D3557", "D4A72C", "17212B", "667085", "F3F6F9", "FBFCFE"),
    OfficeTheme("森林研究", "26734D", "8EBF63", "17342A", "66756E", "EDF7F0", "FAFCFA"),
    OfficeTheme("暖调创意", "B4532A", "E9A23B", "34251F", "756A65", "FFF3E8", "FFFCF8"),
)


def _resolve_theme(title: str, style_instructions: str) -> OfficeTheme:
    hint = f"{title}\n{style_instructions}".lower()
    if any(word in hint for word in ("科技", "tech", "cyber", "数据", "代码", "software")):
        return THEMES[0]
    if any(word in hint for word in ("学术", "academic", "论文", "研究", "formal")):
        return THEMES[1]
    if any(word in hint for word in ("自然", "生态", "green", "biology", "environment")):
        return THEMES[2]
    if any(word in hint for word in ("温暖", "warm", "人文", "创意", "creative")):
        return THEMES[3]
    digest = hashlib.sha256(hint.encode("utf-8")).digest()
    return THEMES[digest[0] % len(THEMES)]


def _rgb(hex_color: str) -> tuple[int, int, int]:
    value = hex_color.lstrip("#")
    return tuple(int(value[index:index + 2], 16) for index in (0, 2, 4))  # type: ignore[return-value]


def _short_keywords(title: str, sections: list[tuple[str, list[str]]], limit: int = 5) -> list[str]:
    candidates = [title, *(name for name, _ in sections)]
    result: list[str] = []
    for candidate in candidates:
        clean = re.sub(r"[《》：:，,。.!！?？（）()\[\]【】]", " ", _plain(candidate))
        for item in clean.split():
            item = item.strip()
            if 2 <= len(item) <= 12 and item not in result:
                result.append(item)
            if len(result) >= limit:
                return result
    return result or ["理解", "连接", "应用"]


def _theme_art(theme: OfficeTheme, seed_text: str, width: int = 1400, height: int = 560) -> BytesIO:
    """Create a deterministic, topic-specific abstract visual without external image services."""
    accent = _rgb(theme.accent)
    accent_2 = _rgb(theme.accent_2)
    canvas = _rgb(theme.canvas)
    digest = hashlib.sha256(seed_text.encode("utf-8")).digest()
    image = Image.new("RGB", (width, height), canvas)
    draw = ImageDraw.Draw(image, "RGBA")

    for y in range(height):
        ratio = y / max(1, height - 1)
        color = tuple(int(canvas[index] * (1 - ratio * 0.16) + accent_2[index] * ratio * 0.16) for index in range(3))
        draw.line((0, y, width, y), fill=(*color, 255))

    points: list[tuple[int, int, int]] = []
    for index in range(10):
        x = 80 + ((digest[index] * 47 + index * 137) % (width - 160))
        y = 55 + ((digest[index + 10] * 31 + index * 83) % (height - 110))
        radius = 16 + digest[index + 20] % 52
        points.append((x, y, radius))
    for index, (x, y, _) in enumerate(points):
        for target in (index + 1, index + 3):
            tx, ty, _ = points[target % len(points)]
            draw.line((x, y, tx, ty), fill=(*accent, 42), width=3)
    for index, (x, y, radius) in enumerate(points):
        color = accent if index % 3 else accent_2
        draw.ellipse((x - radius, y - radius, x + radius, y + radius), fill=(*color, 44), outline=(*color, 125), width=3)
        draw.ellipse((x - 6, y - 6, x + 6, y + 6), fill=(*color, 220))

    output = BytesIO()
    image.save(output, format="PNG", optimize=True)
    output.seek(0)
    return output


def _plain(text: str) -> str:
    text = re.sub(r"!\[([^]]*)]\([^)]+\)", r"\1", text)
    text = re.sub(r"\[([^]]+)]\([^)]+\)", r"\1", text)
    return re.sub(r"[*_`~]", "", text).strip()


def _safe_lines(content: str) -> list[str]:
    return [line.rstrip() for line in content.replace("\r\n", "\n").replace("\r", "\n").split("\n")]


def _presentation_sections(content: str, fallback_title: str) -> list[tuple[str, list[str]]]:
    sections: list[tuple[str, list[str]]] = []
    current_title = "核心内容"
    current_items: list[str] = []

    def flush() -> None:
        nonlocal current_items
        if current_items:
            sections.append((current_title, current_items))
            current_items = []

    for raw in _safe_lines(content):
        line = raw.strip()
        if not line or line.startswith("```"):
            continue
        if line.startswith("# "):
            continue
        if line.startswith("## "):
            flush()
            current_title = _plain(line[3:]) or fallback_title
            continue
        if line.startswith("### "):
            current_items.append(_plain(line[4:]))
            continue
        match = re.match(r"^(?:[-*+]\s+|\d+[.)]\s+)(.+)$", line)
        current_items.append(_plain(match.group(1) if match else line))
    flush()
    return sections or [(fallback_title, ["围绕主题梳理核心概念、关键方法与应用。"])]


def _chunks(items: list[str], size: int = 6) -> Iterable[list[str]]:
    for index in range(0, len(items), size):
        yield items[index:index + size]


def _ppt_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int,
    color: tuple[int, int, int] = (32, 33, 35),
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = anchor
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = PptRGBColor(*color)
    return box


def _ppt_rule(slide, left: float, top: float, width: float, color: tuple[int, int, int] = (184, 188, 196)) -> None:
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        PptInches(left),
        PptInches(top),
        PptInches(width),
        PptInches(0.018),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PptRGBColor(*color)
    line.line.fill.background()


def _ppt_vertical_rule(slide, left: float, top: float, height: float, color: tuple[int, int, int] = (184, 188, 196)) -> None:
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        PptInches(left),
        PptInches(top),
        PptInches(0.018),
        PptInches(height),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PptRGBColor(*color)
    line.line.fill.background()


def _ppt_chrome(slide, page_number: int, section_label: str, theme: OfficeTheme) -> None:
    _ppt_text(slide, f"智学 AI  /  {theme.name}", 0.58, 0.32, 4.2, 0.28, size=10, color=_rgb(theme.accent), bold=True)
    _ppt_text(slide, f"{page_number:02d}", 11.95, 0.31, 0.75, 0.28, size=10, color=(112, 116, 126), align=PP_ALIGN.RIGHT)
    _ppt_rule(slide, 0.58, 0.82, 12.15)
    _ppt_text(slide, section_label[:48], 0.58, 7.02, 8.6, 0.22, size=8, color=(142, 145, 154))


def _ppt_card(slide, left: float, top: float, width: float, height: float, theme: OfficeTheme, *, strong: bool = False):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        PptInches(left),
        PptInches(top),
        PptInches(width),
        PptInches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = PptRGBColor(*_rgb(theme.accent if strong else theme.soft))
    card.line.color.rgb = PptRGBColor(*_rgb(theme.accent if strong else "D9E2E8"))
    card.line.width = PptPt(0.8)
    return card


def _ppt_keyword_visual(slide, keywords: list[str], theme: OfficeTheme, left: float, top: float, width: float, height: float) -> None:
    """Build a topic-specific visual from editable PowerPoint shapes."""
    center_x = left + width / 2
    center_y = top + height / 2
    positions = (
        (left + 0.08, top + 0.18),
        (left + width - 1.55, top + 0.28),
        (left + 0.18, top + height - 0.82),
        (left + width - 1.68, top + height - 0.9),
    )
    center = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        PptInches(center_x - 0.72),
        PptInches(center_y - 0.55),
        PptInches(1.44),
        PptInches(1.1),
    )
    center.fill.solid()
    center.fill.fore_color.rgb = PptRGBColor(*_rgb(theme.accent))
    center.line.fill.background()
    _ppt_text(
        slide,
        keywords[0][:8],
        center_x - 0.56,
        center_y - 0.18,
        1.12,
        0.36,
        size=15,
        color=(255, 255, 255),
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    for index, (x, y) in enumerate(positions):
        label = keywords[(index + 1) % len(keywords)][:8]
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            PptInches(center_x),
            PptInches(center_y),
            PptInches(x + 0.65),
            PptInches(y + 0.28),
        )
        connector.line.color.rgb = PptRGBColor(*_rgb(theme.accent_2))
        connector.line.width = PptPt(1.4)
        node = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            PptInches(x),
            PptInches(y),
            PptInches(1.35),
            PptInches(0.58),
        )
        node.fill.solid()
        node.fill.fore_color.rgb = PptRGBColor(255, 255, 255)
        node.line.color.rgb = PptRGBColor(*_rgb(theme.accent_2))
        _ppt_text(slide, label, x + 0.08, y + 0.16, 1.19, 0.22, size=10, color=_rgb(theme.ink), bold=True, align=PP_ALIGN.CENTER)


def _ppt_bullet_column(slide, items: list[str], left: float, top: float, width: float, height: float, size: int = 19) -> None:
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"—  {item[:150]}"
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = PptPt(size)
        paragraph.font.color.rgb = PptRGBColor(52, 54, 61)
        paragraph.space_after = PptPt(15)
        paragraph.line_spacing = 1.08


def build_pptx(
    title: str,
    subtitle: str,
    content: str,
    *,
    skill_names: list[str] | None = None,
    style_instructions: str = "",
) -> BytesIO:
    deck = Presentation()
    deck.slide_width = PptInches(13.333)
    deck.slide_height = PptInches(7.5)
    blank = deck.slide_layouts[6]
    sections = _presentation_sections(content, title)
    theme = _resolve_theme(title, f"{' '.join(skill_names or [])}\n{style_instructions}")
    keywords = _short_keywords(title, sections)
    accent = _rgb(theme.accent)
    accent_2 = _rgb(theme.accent_2)
    ink = _rgb(theme.ink)
    muted = _rgb(theme.muted)

    cover = deck.slides.add_slide(blank)
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = PptRGBColor(*_rgb(theme.canvas))
    _ppt_text(cover, f"智学 AI  ·  {theme.name}", 0.62, 0.48, 4.5, 0.36, size=13, color=accent, bold=True)
    _ppt_text(
        cover,
        title[:72],
        0.62,
        1.48,
        7.25,
        2.2,
        size=48 if len(title) <= 22 else 40,
        color=ink,
        bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _ppt_rule(cover, 0.62, 4.24, 2.05, accent)
    _ppt_text(cover, subtitle[:140], 0.62, 4.62, 7.1, 0.72, size=20, color=muted)
    _ppt_text(cover, "理解  ·  连接  ·  应用", 0.62, 6.72, 4.8, 0.28, size=11, color=muted)
    _ppt_card(cover, 8.35, 1.18, 4.35, 4.85, theme)
    cover.shapes.add_picture(
        _theme_art(theme, title),
        PptInches(8.46),
        PptInches(1.3),
        width=PptInches(4.13),
        height=PptInches(4.61),
    )
    _ppt_keyword_visual(cover, keywords, theme, 8.68, 1.68, 3.7, 3.7)
    _ppt_text(cover, "01", 11.45, 6.32, 1.25, 0.42, size=20, color=accent_2, bold=True, align=PP_ALIGN.RIGHT)

    page_number = 2
    overview = deck.slides.add_slide(blank)
    overview.background.fill.solid()
    overview.background.fill.fore_color.rgb = PptRGBColor(255, 255, 255)
    _ppt_chrome(overview, page_number, title, theme)
    _ppt_text(overview, "这次讲解将解决什么", 0.62, 1.18, 8.8, 0.7, size=35, color=ink, bold=True)
    _ppt_text(overview, "沿着问题逐步建立理解，最后落实到应用。", 0.62, 1.95, 7.4, 0.42, size=16, color=muted)
    route_titles = [section_title for section_title, _ in sections][:10]
    split_at = (len(route_titles) + 1) // 2
    for column, column_items in enumerate((route_titles[:split_at], route_titles[split_at:])):
        left = 0.68 + column * 6.05
        for index, item in enumerate(column_items):
            top = 2.72 + index * 0.72
            number = index + 1 + (split_at if column else 0)
            _ppt_text(overview, f"{number:02d}", left, top, 0.48, 0.3, size=11, color=accent, bold=True)
            _ppt_text(overview, item[:34], left + 0.62, top - 0.03, 5.12, 0.42, size=17, bold=True)
            _ppt_rule(overview, left + 0.62, top + 0.47, 4.92, (232, 232, 236))

    content_slide_index = 0
    for section_title, items in sections:
        for part_index, part in enumerate(_chunks(items)):
            page_number += 1
            content_slide_index += 1
            slide = deck.slides.add_slide(blank)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = PptRGBColor(255, 255, 255)
            suffix = "（续）" if part_index else ""
            display_title = f"{section_title}{suffix}"[:54]
            _ppt_chrome(slide, page_number, title, theme)
            _ppt_text(slide, display_title, 0.62, 1.12, 11.6, 0.78, size=34 if len(display_title) <= 24 else 29, color=ink, bold=True)

            layout = content_slide_index % 4
            if layout == 1 and len(part) >= 2:
                _ppt_card(slide, 0.62, 2.18, 5.55, 3.83, theme)
                _ppt_text(slide, part[0][:150], 0.92, 2.5, 4.95, 2.7, size=27, color=ink, bold=True, anchor=MSO_ANCHOR.MIDDLE)
                _ppt_vertical_rule(slide, 6.31, 2.24, 3.72, accent)
                _ppt_bullet_column(slide, part[1:], 6.72, 2.22, 5.85, 3.7, 18)
            elif layout == 2 and len(part) >= 4:
                midpoint = (len(part) + 1) // 2
                _ppt_card(slide, 0.62, 2.12, 5.62, 4.0, theme)
                _ppt_card(slide, 6.66, 2.12, 5.98, 4.0, theme)
                _ppt_text(slide, "先理解", 0.9, 2.42, 5.0, 0.36, size=14, color=accent, bold=True)
                _ppt_bullet_column(slide, part[:midpoint], 0.68, 2.82, 5.45, 3.25, 18)
                _ppt_text(slide, "再应用", 6.94, 2.42, 5.3, 0.36, size=14, color=accent, bold=True)
                _ppt_bullet_column(slide, part[midpoint:], 6.78, 2.82, 5.45, 3.25, 18)
            elif layout == 3 and len(part) >= 3:
                for index, item in enumerate(part[:3]):
                    left = 0.68 + index * 4.08
                    _ppt_card(slide, left, 2.18, 3.72, 3.55, theme, strong=index == 0)
                    text_color = (255, 255, 255) if index == 0 else ink
                    _ppt_text(slide, f"0{index + 1}", left + 0.2, 2.43, 0.7, 0.35, size=13, color=(255, 255, 255) if index == 0 else accent, bold=True)
                    _ppt_text(slide, item[:125], left + 0.2, 3.02, 3.25, 2.25, size=19, color=text_color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
                if len(part) > 3:
                    _ppt_text(slide, "补充：" + "；".join(part[3:])[:150], 0.68, 5.76, 11.8, 0.52, size=14, color=(100, 104, 113))
            else:
                first = part[0] if part else section_title
                _ppt_text(slide, first[:165], 0.68, 2.22, 7.1, 1.62, size=29, color=ink, bold=True, anchor=MSO_ANCHOR.MIDDLE)
                _ppt_keyword_visual(slide, [section_title, *keywords], theme, 8.18, 2.05, 4.22, 3.72)
                _ppt_rule(slide, 0.68, 4.18, 6.9, accent)
                _ppt_bullet_column(slide, part[1:] or ["用自己的语言复述这一点，并尝试连接到一个具体问题。"], 0.68, 4.62, 6.95, 1.55, 18)

    page_number += 1
    close = deck.slides.add_slide(blank)
    close.background.fill.solid()
    close.background.fill.fore_color.rgb = PptRGBColor(*_rgb(theme.canvas))
    last_items = sections[-1][1][:3] if sections else []
    _ppt_text(close, f"智学 AI  ·  {theme.name}", 0.62, 0.48, 4.5, 0.34, size=12, color=accent, bold=True)
    _ppt_text(close, "现在，把理解变成行动", 0.62, 1.58, 10.8, 1.2, size=43, color=ink, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _ppt_text(close, "回到题目或实践场景，完成一次独立解释与应用。", 0.62, 3.02, 9.6, 0.5, size=19, color=muted)
    for index, item in enumerate(last_items or ["复述核心概念", "完成一道变式题", "记录并订正易错点"]):
        _ppt_text(close, f"{index + 1}", 0.68 + index * 4.04, 4.42, 0.35, 0.36, size=13, color=accent, bold=True)
        _ppt_text(close, item[:72], 1.12 + index * 4.04, 4.36, 3.25, 1.15, size=17, bold=True)
    _ppt_text(close, f"{page_number:02d}", 11.45, 6.34, 1.25, 0.7, size=32, color=accent_2, bold=True, align=PP_ALIGN.RIGHT)

    output = BytesIO()
    deck.save(output)
    output.seek(0)
    return output


def _set_run_font(run, name: str, size: int, color: str = "202123", bold: bool = False) -> None:
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    run._element.get_or_add_rPr().rFonts.set(qn("w:eastAsia"), name)


def _shade_paragraph(paragraph, fill: str) -> None:
    properties = paragraph._p.get_or_add_pPr()
    shading = OxmlElement("w:shd")
    shading.set(qn("w:fill"), fill)
    properties.append(shading)


def _shade_cell(cell, fill: str) -> None:
    properties = cell._tc.get_or_add_tcPr()
    shading = properties.find(qn("w:shd"))
    if shading is None:
        shading = OxmlElement("w:shd")
        properties.append(shading)
    shading.set(qn("w:fill"), fill)


def _set_cell_margins(cell, top: int = 130, start: int = 150, bottom: int = 130, end: int = 150) -> None:
    properties = cell._tc.get_or_add_tcPr()
    margins = properties.first_child_found_in("w:tcMar")
    if margins is None:
        margins = OxmlElement("w:tcMar")
        properties.append(margins)
    for edge, value in (("top", top), ("start", start), ("bottom", bottom), ("end", end)):
        node = margins.find(qn(f"w:{edge}"))
        if node is None:
            node = OxmlElement(f"w:{edge}")
            margins.append(node)
        node.set(qn("w:w"), str(value))
        node.set(qn("w:type"), "dxa")


def _add_doc_visual_map(document: Document, keywords: list[str], theme: OfficeTheme) -> None:
    label = document.add_paragraph()
    label.paragraph_format.space_before = Pt(12)
    label.paragraph_format.space_after = Pt(7)
    run = label.add_run("内容地图")
    _set_run_font(run, "Microsoft YaHei", 10, theme.accent, True)

    table = document.add_table(rows=1, cols=min(4, max(2, len(keywords))))
    table.autofit = False
    for index, cell in enumerate(table.rows[0].cells):
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        _shade_cell(cell, theme.soft if index else theme.accent)
        _set_cell_margins(cell, 190, 150, 190, 150)
        paragraph = cell.paragraphs[0]
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = paragraph.add_run(keywords[index % len(keywords)][:12])
        _set_run_font(
            run,
            "Microsoft YaHei",
            10,
            "FFFFFF" if index == 0 else theme.ink,
            True,
        )
    after = document.add_paragraph()
    after.paragraph_format.space_after = Pt(8)


def build_docx(
    title: str,
    subtitle: str,
    content: str,
    *,
    skill_names: list[str] | None = None,
    style_instructions: str = "",
) -> BytesIO:
    sections = _presentation_sections(content, title)
    theme = _resolve_theme(title, f"{' '.join(skill_names or [])}\n{style_instructions}")
    keywords = _short_keywords(title, sections)
    document = Document()
    section = document.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(0.9)
    section.right_margin = Inches(1)
    section.bottom_margin = Inches(0.9)
    section.left_margin = Inches(1)
    section.header_distance = Inches(0.45)
    section.footer_distance = Inches(0.45)

    normal = document.styles["Normal"]
    normal.font.name = "Microsoft YaHei"
    normal.font.size = Pt(10.5)
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    normal.paragraph_format.space_after = Pt(6)
    normal.paragraph_format.line_spacing = 1.2
    for style_name, size, color, before, after in (
        ("Heading 1", 18, theme.accent, 16, 8),
        ("Heading 2", 14, theme.accent, 12, 6),
        ("Heading 3", 12, theme.ink, 9, 4),
    ):
        style = document.styles[style_name]
        style.font.name = "Microsoft YaHei"
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = RGBColor.from_string(color)
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
        style.paragraph_format.space_before = Pt(before)
        style.paragraph_format.space_after = Pt(after)
        style.paragraph_format.keep_with_next = True

    title_paragraph = document.add_paragraph()
    title_paragraph.paragraph_format.space_after = Pt(7)
    title_run = title_paragraph.add_run(title[:120])
    _set_run_font(title_run, "Microsoft YaHei", 26, theme.ink, True)
    subtitle_paragraph = document.add_paragraph()
    subtitle_paragraph.paragraph_format.space_after = Pt(5)
    subtitle_run = subtitle_paragraph.add_run(subtitle[:180])
    _set_run_font(subtitle_run, "Microsoft YaHei", 10, theme.muted)
    theme_paragraph = document.add_paragraph()
    theme_paragraph.paragraph_format.space_after = Pt(4)
    theme_run = theme_paragraph.add_run(
        f"{theme.name}  ·  {(' / '.join(skill_names or [])) if skill_names else '自适应视觉编排'}"
    )
    _set_run_font(theme_run, "Microsoft YaHei", 8, theme.accent, True)
    _shade_paragraph(theme_paragraph, theme.soft)
    art_paragraph = document.add_paragraph()
    art_paragraph.paragraph_format.space_before = Pt(5)
    art_paragraph.paragraph_format.space_after = Pt(4)
    art_paragraph.add_run().add_picture(_theme_art(theme, title), width=Inches(6.45))
    _add_doc_visual_map(document, keywords, theme)

    in_code = False
    code_lines: list[str] = []
    for raw in _safe_lines(content):
        line = raw.rstrip()
        if line.strip().startswith("```"):
            if in_code:
                paragraph = document.add_paragraph()
                paragraph.paragraph_format.left_indent = Inches(0.18)
                paragraph.paragraph_format.right_indent = Inches(0.18)
                paragraph.paragraph_format.space_before = Pt(4)
                paragraph.paragraph_format.space_after = Pt(9)
                _shade_paragraph(paragraph, theme.soft)
                run = paragraph.add_run("\n".join(code_lines))
                _set_run_font(run, "Consolas", 9, theme.ink)
                code_lines = []
                in_code = False
            else:
                in_code = True
            continue
        if in_code:
            code_lines.append(line)
            continue
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("# "):
            continue
        if stripped.startswith("### "):
            document.add_paragraph(_plain(stripped[4:]), style="Heading 3")
        elif stripped.startswith("## "):
            document.add_paragraph(_plain(stripped[3:]), style="Heading 2")
        elif re.match(r"^[-*+]\s+", stripped):
            paragraph = document.add_paragraph(style="List Bullet")
            paragraph.paragraph_format.space_after = Pt(4)
            paragraph.add_run(_plain(re.sub(r"^[-*+]\s+", "", stripped)))
        elif re.match(r"^\d+[.)]\s+", stripped):
            paragraph = document.add_paragraph(style="List Number")
            paragraph.paragraph_format.space_after = Pt(4)
            paragraph.add_run(_plain(re.sub(r"^\d+[.)]\s+", "", stripped)))
        else:
            paragraph = document.add_paragraph()
            paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
            paragraph.add_run(_plain(stripped))

    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    footer_run = footer.add_run(subtitle[:70])
    _set_run_font(footer_run, "Microsoft YaHei", 8, theme.muted)

    output = BytesIO()
    document.save(output)
    output.seek(0)
    return output
